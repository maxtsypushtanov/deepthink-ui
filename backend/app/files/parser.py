"""Stage 1 — Parse: extract text from uploaded files."""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".txt", ".md", ".py", ".js", ".ts", ".tsx", ".jsx", ".json",
                        ".yaml", ".yml", ".toml", ".csv", ".html", ".css", ".sql",
                        ".sh", ".bash", ".rs", ".go", ".java", ".c", ".cpp", ".h",
                        ".rb", ".php", ".swift", ".kt", ".xml", ".ini", ".cfg",
                        ".pdf", ".docx", ".pptx", ".xlsx", ".xls",
                        ".png", ".jpg", ".jpeg", ".gif", ".webp"}

# Code-like extensions for syntax hint
CODE_EXTENSIONS = {".py", ".js", ".ts", ".tsx", ".jsx", ".rs", ".go", ".java",
                   ".c", ".cpp", ".h", ".rb", ".php", ".swift", ".kt", ".sql",
                   ".sh", ".bash", ".css", ".html"}


def parse_file(path: str | Path) -> dict:
    """Extract text and metadata from a file.

    Returns:
        {
            "text": str,
            "filename": str,
            "extension": str,
            "char_count": int,
            "file_type": "text" | "code" | "pdf" | "docx",
            "error": str | None,
        }
    """
    p = Path(path)
    ext = p.suffix.lower()
    filename = p.name
    result = {
        "filename": filename,
        "extension": ext,
        "text": "",
        "char_count": 0,
        "file_type": "text",
        "error": None,
    }

    if ext not in SUPPORTED_EXTENSIONS:
        result["error"] = f"Неподдерживаемый формат: {ext}"
        return result

    try:
        if ext == ".pdf":
            result["text"] = _parse_pdf(p)
            result["file_type"] = "pdf"
        elif ext == ".docx":
            result["text"] = _parse_docx(p)
            result["file_type"] = "docx"
        elif ext == ".pptx":
            result["text"] = _parse_pptx(p)
            result["file_type"] = "pptx"
        elif ext in (".xlsx", ".xls"):
            result["text"] = _parse_xlsx(p)
            result["file_type"] = "xlsx"
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            result["file_type"] = "image"
            result["image_base64"] = _image_to_base64(p)
            result["image_mime"] = f"image/{'jpeg' if ext in ('.jpg', '.jpeg') else ext.lstrip('.')}"
            result["text"] = f"[Изображение: {p.name}]"
        elif ext in CODE_EXTENSIONS:
            result["text"] = p.read_text(encoding="utf-8", errors="replace")
            result["file_type"] = "code"
        else:
            result["text"] = p.read_text(encoding="utf-8", errors="replace")
            result["file_type"] = "text"
    except Exception as e:
        logger.warning("Failed to parse %s: %s", filename, e)
        result["error"] = str(e)

    result["char_count"] = len(result["text"])

    # TRIZ #6: Structural pre-scan — extract metadata for smart file analysis
    result["structure"] = _quick_scan(result["text"], result["file_type"], ext)

    return result


def _quick_scan(text: str, file_type: str, ext: str) -> dict:
    """Fast structural analysis without LLM. Returns metadata for smart prompts."""
    import re as _re
    info: dict = {}

    if file_type == "code":
        # Extract function/class names
        functions = _re.findall(r'(?:def|function|func|fn|async\s+def|async\s+function)\s+(\w+)', text)
        classes = _re.findall(r'(?:class|struct|interface|enum|trait|impl)\s+(\w+)', text)
        lines = text.count('\n') + 1
        info["functions"] = functions[:20]
        info["classes"] = classes[:10]
        info["lines"] = lines
        info["summary"] = f"{len(functions)} функций, {len(classes)} классов, {lines} строк"

    elif file_type == "pdf":
        pages = text.count('[Страница ')
        # Extract section headings (lines in ALL CAPS or starting with digits)
        headings = _re.findall(r'^\s*(?:\d+[\.\)]\s+)?[A-ZА-ЯЁ][A-ZА-ЯЁ\s]{5,}$', text, _re.MULTILINE)
        info["pages"] = pages
        info["headings"] = [h.strip() for h in headings[:10]]
        info["summary"] = f"{pages} стр., {len(headings)} секций"

    elif file_type in ("xlsx", "xls"):
        sheets = _re.findall(r'\[Лист: (.+?)\]', text)
        rows = text.count('\n')
        info["sheets"] = sheets
        info["rows"] = rows
        info["summary"] = f"{len(sheets)} листов, ~{rows} строк"

    elif file_type == "docx":
        paragraphs = len([p for p in text.split('\n\n') if p.strip()])
        words = len(text.split())
        info["paragraphs"] = paragraphs
        info["words"] = words
        info["summary"] = f"{paragraphs} абзацев, {words} слов"

    elif file_type == "pptx":
        slides = text.count('[Слайд ')
        info["slides"] = slides
        info["summary"] = f"{slides} слайдов"

    elif ext in ('.csv',):
        lines_list = text.split('\n')
        if lines_list:
            columns = lines_list[0].split(',') if ',' in lines_list[0] else lines_list[0].split('\t')
            info["columns"] = [c.strip().strip('"') for c in columns[:20]]
            info["rows"] = len(lines_list) - 1
            info["summary"] = f"{len(columns)} столбцов, {len(lines_list)-1} строк"

    else:
        lines = text.count('\n') + 1
        words = len(text.split())
        info["lines"] = lines
        info["words"] = words
        info["summary"] = f"{words} слов, {lines} строк"

    return info


def _parse_pdf(path: Path) -> str:
    """Extract text from PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    total_pages = len(reader.pages)
    pages = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"[Страница {i + 1}]\n{text}")

    if not pages:
        # No text extracted — likely scanned/image PDF
        raise ValueError(
            f"PDF содержит {total_pages} стр., но текст не извлекается. "
            "Возможно, это сканированный документ без текстового слоя."
        )

    return "\n\n".join(pages)


def _parse_pptx(path: Path) -> str:
    """Extract text from PowerPoint PPTX."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"[Слайд {i + 1}]\n" + "\n".join(texts))

    if not slides:
        raise ValueError(
            f"Презентация содержит {len(prs.slides)} слайдов, но текст не найден."
        )
    return "\n\n".join(slides)


def _parse_xlsx(path: Path) -> str:
    """Extract data from Excel XLSX/XLS as text tables."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            header = f"[Лист: {ws.title}]"
            sheets.append(header + "\n" + "\n".join(rows))
    wb.close()

    if not sheets:
        raise ValueError("Excel-файл пуст или не содержит данных.")
    return "\n\n".join(sheets)


def _parse_image(path: Path) -> str:
    """Return a placeholder description for images (actual analysis done by vision LLM)."""
    size_kb = path.stat().st_size / 1024
    return f"[Изображение: {path.name}, {size_kb:.0f} КБ]"


def _image_to_base64(path: Path) -> str:
    """Encode image file as base64 string."""
    import base64
    return base64.b64encode(path.read_bytes()).decode("ascii")


def _parse_docx(path: Path) -> str:
    """Extract text from DOCX preserving heading structure for Cognitive Map.

    Headings are converted to markdown-style markers so the section splitter
    can detect them: # Heading 1, ## Heading 2, ### Heading 3, etc.
    """
    from docx import Document

    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name or "").lower()
        if "heading 1" in style or "заголовок 1" in style:
            parts.append(f"# {text}")
        elif "heading 2" in style or "заголовок 2" in style:
            parts.append(f"## {text}")
        elif "heading 3" in style or "заголовок 3" in style:
            parts.append(f"### {text}")
        elif "heading 4" in style or "заголовок 4" in style:
            parts.append(f"#### {text}")
        elif "heading" in style or "заголовок" in style:
            parts.append(f"## {text}")
        elif "title" in style or "название" in style:
            parts.append(f"# {text}")
        elif "subtitle" in style or "подзаголовок" in style:
            parts.append(f"## {text}")
        else:
            parts.append(text)

    # Also extract tables if present
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"\n### [Таблица {i + 1}]")
            parts.append("\n".join(rows))

    return "\n\n".join(parts)
